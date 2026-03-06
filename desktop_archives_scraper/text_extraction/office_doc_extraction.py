"""Office document text extraction: Word, PowerPoint, and Excel.

Extraction strategy per format
-------------------------------
Modern formats (.docx/.pptx/.xlsx) are handled in-process with pure-Python
libraries (mammoth, python-docx, python-pptx, openpyxl) unless the file
exceeds ``OFFICE_SUBPROCESS_FILESIZE_THRESHOLD_BYTES``, in which case they are
forwarded to a child worker process to protect the main worker from OOM.

Legacy formats (.doc/.ppt/.pps/.xls) are *always* handled in a child worker
process.  Inside that worker the extraction chain is:

1. **Windows COM** (``Word.Application`` / ``PowerPoint.Application`` /
   ``Excel.Application``) — converts to the modern XML format in a temp dir,
   then parses with the same in-process library.  Requires MS Office to be
   installed; uses ``pywin32``.
2. **LibreOffice** (``soffice --headless --convert-to``) — fallback when COM
   is unavailable or fails.  Requires LibreOffice to be installed and on PATH
   (or pointed to via the ``soffice_path`` config key).

Text normalisation
------------------
All extracted text passes through ``_normalize_office_text()`` which strips
null bytes and non-printable control characters, normalises line endings, and
collapses excessive blank lines before returning.
"""
import logging
import os
import re
import json
import shutil
import shlex
import subprocess
import sys
import tempfile
import time
import uuid

from pathlib import Path
from typing import List

from .basic_extraction import FileTextExtractor, TextExtractionError
from .extraction_utils import validate_file

logger = logging.getLogger(__name__)

# Files larger than this threshold are sent to a child worker process even for
# modern formats, to avoid OOM in the main worker.
OFFICE_SUBPROCESS_FILESIZE_THRESHOLD_BYTES = int(
    os.getenv("OFFICE_SUBPROCESS_FILESIZE_THRESHOLD_BYTES", str(75 * 1024 * 1024))
)
# Default wall-clock timeout for the child worker subprocess (seconds).
OFFICE_WORKER_TIMEOUT_S = int(os.getenv("OFFICE_WORKER_TIMEOUT_S", "180"))
# Env var name set to "1" inside the worker subprocess to prevent infinite recursion.
OFFICE_WORKER_MODE_ENV = "OFFICE_WORKER_MODE"
# Max chars of subprocess stderr to include in error messages.
OFFICE_WORKER_STDERR_TAIL_MAX_CHARS = 4000
# Python module path used to spawn the worker subprocess.
OFFICE_WORKER_MODULE = "desktop_archives_scraper.text_extraction.office_extraction_worker"

_OFFICE_WORKER_MEM_MB_RAW = os.getenv("OFFICE_WORKER_MEM_MB")
if _OFFICE_WORKER_MEM_MB_RAW is None:
    OFFICE_WORKER_MEM_MB = 2048
else:
    try:
        OFFICE_WORKER_MEM_MB = int(_OFFICE_WORKER_MEM_MB_RAW)
    except ValueError:
        logger.warning("Invalid OFFICE_WORKER_MEM_MB=%s; defaulting to 2048", _OFFICE_WORKER_MEM_MB_RAW)
        OFFICE_WORKER_MEM_MB = 2048

# Legacy formats are always sent to a subprocess regardless of file size.
OFFICE_LEGACY_ALWAYS_SUBPROCESS = {"doc", "ppt", "pps", "xls"}
# Modern formats are only sent to a subprocess when they exceed the size threshold.
OFFICE_MODERN_SIZE_GATED_SUBPROCESS = {"docx", "docm", "pptx", "pptm", "ppsx", "xlsx", "xlsm"}


def _env_flag(name: str, default: bool = False) -> bool:
    """Read an environment variable as a boolean flag.

    Truthy string values: ``"1"``, ``"true"``, ``"yes"``, ``"on"``
    (case-insensitive).  Returns *default* when the variable is unset.
    """
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def _tail_text(text: str | None, max_chars: int = OFFICE_WORKER_STDERR_TAIL_MAX_CHARS) -> str:
    """Return the last *max_chars* characters of *text*, stripped of whitespace.

    Used to include a bounded tail of subprocess stderr in error messages
    without risking unbounded log growth for very verbose Office tools.
    """
    if not text:
        return ""
    cleaned = text.strip()
    if len(cleaned) <= max_chars:
        return cleaned
    return f"...[truncated] {cleaned[-max_chars:]}"


def _in_office_worker_mode() -> bool:
    """Return True when running inside the isolated office worker subprocess.

    The parent sets ``OFFICE_WORKER_MODE=1`` before spawning the worker so
    that extractor ``__call__`` methods bypass the subprocess-routing guard
    and perform extraction directly in-process.
    """
    return _env_flag(OFFICE_WORKER_MODE_ENV, default=False)


def _should_route_to_office_subprocess(source: Path, ext: str) -> bool:
    """Decide whether this file should be handled by a child worker process.

    Returns False immediately when already running inside a worker (prevents
    infinite recursion).  Legacy formats always go to a subprocess.  Modern
    formats only do so when the file exceeds the size threshold.
    """
    if _in_office_worker_mode():
        return False

    if ext in OFFICE_LEGACY_ALWAYS_SUBPROCESS:
        return True

    if ext in OFFICE_MODERN_SIZE_GATED_SUBPROCESS:
        return source.stat().st_size >= OFFICE_SUBPROCESS_FILESIZE_THRESHOLD_BYTES

    return False


def run_office_worker(
    input_path: Path | str,
    *,
    timeout_s: int = OFFICE_WORKER_TIMEOUT_S,
    mem_mb: int | None = OFFICE_WORKER_MEM_MB,
    config: dict | None = None,
    worker_cmd: list[str] | None = None,
) -> str:
    """Spawn the office extraction worker subprocess and return the extracted text.

    The worker is invoked as::

        python -m desktop_archives_scraper.text_extraction.office_extraction_worker \
            --input <path> --config-json <json> [--timeout-s N] [--mem-mb N]

    The worker writes a single JSON object to stdout.  This function parses
    that payload and returns the ``"text"`` field on success, or raises
    ``TextExtractionError`` on any failure.

    Parameters
    ----------
    input_path:
        Path to the Office document to process.
    timeout_s:
        Wall-clock seconds before the subprocess is killed.
    mem_mb:
        Soft memory cap passed to the worker (Linux only; ignored on Windows).
    config:
        Extractor config dict forwarded to the worker as JSON.
    worker_cmd:
        Override the default ``[sys.executable, "-m", MODULE]`` command.
        Useful for testing.

    Raises
    ------
    TextExtractionError
        On subprocess timeout, crash (empty stdout), JSON parse failure, or
        when the worker reports ``ok=False``.
    """
    source = validate_file(str(input_path))
    worker_config = config or {}

    cmd_override = os.getenv("OFFICE_WORKER_CMD")
    if worker_cmd is not None:
        cmd = list(worker_cmd)
    elif cmd_override:
        cmd = shlex.split(cmd_override)
    else:
        cmd = [sys.executable, "-m", OFFICE_WORKER_MODULE]

    cmd.extend(["--input", str(source), "--config-json", json.dumps(worker_config)])

    if timeout_s is not None:
        cmd.extend(["--timeout-s", str(timeout_s)])

    if mem_mb is not None:
        cmd.extend(["--mem-mb", str(mem_mb)])

    try:
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=timeout_s,
            check=False,
        )
    except subprocess.TimeoutExpired as exc:
        stderr_tail = _tail_text(exc.stderr)
        raise TextExtractionError(
            f"office worker timed out after {timeout_s}s: reason=timeout retryable=True stderr_tail={stderr_tail}"
        ) from exc

    stderr_tail = _tail_text(result.stderr)
    stdout = (result.stdout or "").strip()
    if not stdout:
        raise TextExtractionError(
            f"office worker crash: reason=worker_crash retryable=True worker_exit_code={result.returncode} worker_stderr_tail={stderr_tail}"
        )

    try:
        payload = json.loads(stdout)
    except json.JSONDecodeError as exc:
        raise TextExtractionError(
            f"office worker crash: reason=worker_crash retryable=True worker_exit_code={result.returncode} worker_stderr_tail={stderr_tail}"
        ) from exc

    if not isinstance(payload, dict):
        raise TextExtractionError(
            f"office worker crash: reason=worker_crash retryable=True worker_exit_code={result.returncode} worker_stderr_tail={stderr_tail}"
        )

    ok = bool(payload.get("ok", False))
    if not ok:
        error = payload.get("error") if isinstance(payload.get("error"), dict) else {}
        reason = str(error.get("reason", "parse_failed"))
        retryable = bool(error.get("retryable", False))
        details = error.get("details") if isinstance(error.get("details"), dict) else {}
        raise TextExtractionError(
            "office worker failed: "
            f"reason={reason} retryable={retryable} details={details} "
            f"worker_exit_code={result.returncode} worker_stderr_tail={stderr_tail}"
        )

    if result.returncode != 0:
        raise TextExtractionError(
            f"office worker crash: reason=worker_crash retryable=True worker_exit_code={result.returncode} worker_stderr_tail={stderr_tail}"
        )

    text = payload.get("text", "")
    if not isinstance(text, str):
        raise TextExtractionError(
            f"office worker crash: reason=worker_crash retryable=True worker_exit_code={result.returncode} worker_stderr_tail={stderr_tail}"
        )

    return text


class OfficeConverter:
    """Convert legacy Office documents to modern formats using headless LibreOffice.

    This is the *fallback* converter used when Windows COM automation is
    unavailable or fails.  Each ``convert()`` call runs LibreOffice in a
    private ``HOME`` directory (via ``--env HOME=...``) to avoid profile
    locking errors when multiple workers run concurrently.
    """

    def __init__(self, soffice_path: str = "soffice", default_timeout_s: int = 90):
        """Initialise the converter.

        Parameters
        ----------
        soffice_path:
            Name or full path of the LibreOffice executable.  Defaults to
            ``"soffice"`` which is resolved via PATH.
        default_timeout_s:
            Default conversion timeout in seconds used when ``convert()`` is
            called without an explicit ``timeout_s``.
        """
        self.soffice_path = soffice_path
        self.default_timeout_s = default_timeout_s

    def convert(self, input_path: Path, target_ext: str, out_dir: Path, *, timeout_s: int | None = None) -> Path:
        """Convert *input_path* to *target_ext* format and place the result in *out_dir*.

        LibreOffice is run with a per-job isolated ``HOME`` directory under
        ``/tmp/archives_scraper/office/<uuid>/`` to prevent profile collisions
        when multiple worker processes run concurrently.  The temp directory is
        cleaned up in the ``finally`` block unless ``DEBUG_KEEP_TEMPS=1``.

        Parameters
        ----------
        input_path:
            Source file to convert.
        target_ext:
            Target format extension without a leading dot, e.g. ``"docx"``.
        out_dir:
            Directory where the converted file is staged after conversion.
        timeout_s:
            Override the instance default timeout.

        Returns
        -------
        Path
            Path to the converted file inside *out_dir*.

        Raises
        ------
        TextExtractionError
            On non-zero LibreOffice exit code, missing/empty output file,
            timeout, or LibreOffice executable not found.
        """
        source = validate_file(str(input_path))
        target = target_ext.lower().lstrip(".")
        timeout = timeout_s if timeout_s is not None else self.default_timeout_s

        out_dir.mkdir(parents=True, exist_ok=True)
        job_id = uuid.uuid4().hex
        job_root = Path("/tmp/archives_scraper/office") / job_id
        lo_home = job_root / "home"
        convert_out = job_root / "out"
        lo_home.mkdir(parents=True, exist_ok=True)
        convert_out.mkdir(parents=True, exist_ok=True)

        cmd = [
            self.soffice_path,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--norestore",
            "--convert-to",
            target,
            "--outdir",
            str(convert_out),
            str(source),
        ]

        start = time.time()
        logger.info(
            "Starting LibreOffice conversion",
            extra={
                "source": str(source),
                "target_ext": target,
                "timeout_s": timeout,
            },
        )

        try:
            env = dict(os.environ)
            env["HOME"] = str(lo_home)

            proc = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                env=env,
                timeout=timeout,
                check=False,
                text=True,
            )

            duration_ms = int((time.time() - start) * 1000)
            if proc.returncode != 0:
                stderr = (proc.stderr or "").strip()
                raise TextExtractionError(
                    f"libreoffice conversion failed: rc={proc.returncode}, target={target}, stderr={stderr[:400]}"
                )

            selected = self._select_output(convert_out, source, target)
            if selected.stat().st_size <= 0:
                raise TextExtractionError(f"libreoffice produced empty output: {selected}")

            output_path = out_dir / f"{source.stem}.{target}"
            shutil.copy2(selected, output_path)
            if not output_path.exists() or output_path.stat().st_size <= 0:
                raise TextExtractionError(f"failed to stage converted output: {output_path}")

            logger.info(
                "Completed LibreOffice conversion",
                extra={
                    "source": str(source),
                    "converted_to": str(output_path),
                    "target_ext": target,
                    "convert_ms": duration_ms,
                },
            )
            return output_path

        except subprocess.TimeoutExpired as exc:
            raise TextExtractionError(
                f"libreoffice conversion timed out after {timeout}s for {source}"
            ) from exc
        except FileNotFoundError as exc:
            raise TextExtractionError(
                f"LibreOffice executable not found: {self.soffice_path}"
            ) from exc
        finally:
            if _env_flag("DEBUG_KEEP_TEMPS", default=False):
                logger.warning("Preserving Office temp directory for debug: %s", job_root)
            else:
                shutil.rmtree(job_root, ignore_errors=True)

    @staticmethod
    def _select_output(convert_out: Path, source: Path, target_ext: str) -> Path:
        """Find the converted output file in *convert_out*.

        LibreOffice uses the source filename stem for the output, but the
        capitalisation may differ.  This method tries an exact match first,
        then a case-insensitive stem match, then accepts the only file if
        there is exactly one.  Raises ``TextExtractionError`` on ambiguity or
        when no output is found.
        """
        exact = convert_out / f"{source.stem}.{target_ext}"
        if exact.exists() and exact.is_file():
            return exact

        target_files = [
            candidate
            for candidate in convert_out.iterdir()
            if candidate.is_file() and candidate.suffix.lower().lstrip(".") == target_ext
        ]
        if not target_files:
            raise TextExtractionError(
                f"libreoffice conversion output not found for {source} in {convert_out}"
            )

        basename_matches = [f for f in target_files if f.stem.casefold() == source.stem.casefold()]
        if len(basename_matches) == 1:
            return basename_matches[0]

        if len(target_files) == 1:
            return target_files[0]

        names = ", ".join(sorted(f.name for f in target_files))
        raise TextExtractionError(
            f"multiple libreoffice outputs for {source.name}: {names}"
        )


def _strip_control_chars(text: str) -> str:
    """Remove non-printable control characters except tab and newline."""
    return "".join(ch for ch in text if ch == "\n" or ch == "\t" or ord(ch) >= 32)


def _normalize_office_text(text: str) -> str:
    """Normalise whitespace and encoding artefacts in extracted Office text.

    - Converts all line endings to ``\n``.
    - Strips null bytes and non-printable control characters.
    - Trims trailing whitespace from each line.
    - Collapses runs of 3+ blank lines to a single blank line.
    - Strips leading/trailing whitespace from the result.
    """
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\x00", "")
    text = _strip_control_chars(text)
    text = "\n".join(line.rstrip() for line in text.split("\n"))
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _truncate_text(text: str, max_chars: int) -> tuple[str, bool]:
    """Truncate *text* to *max_chars* characters.

    Returns ``(text, truncated)`` where *truncated* is True when the input
    exceeded the limit.
    """
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


class WordFileTextExtractor(FileTextExtractor):
    """Extract text from Word documents (.docx, .docm, .doc).

    Extraction chain
    ----------------
    * ``.docx`` / ``.docm``: mammoth (primary) → python-docx (fallback).
    * ``.doc``: Windows COM ``Word.Application`` → ``.docx`` (primary);
      LibreOffice ``soffice --convert-to docx`` (fallback).

    Large or legacy files are forwarded to a child worker subprocess before
    in-process extraction is attempted (see ``_should_route_to_office_subprocess``).
    """

    file_extensions: List[str] = ["docx", "docm", "doc"]

    def __init__(self, converter: OfficeConverter | None = None, max_output_chars: int = 5_000_000):
        """Initialise the extractor.

        Parameters
        ----------
        converter:
            ``OfficeConverter`` instance used as the LibreOffice fallback for
            ``.doc`` files.  A default instance (``soffice`` on PATH) is
            created when not provided.
        max_output_chars:
            Extracted text is truncated to this many characters before
            returning.  Prevents unbounded memory use with very large documents.
        """
        super().__init__()
        self.converter = converter or OfficeConverter()
        self.max_output_chars = max_output_chars
    def __call__(self, path: str) -> str:
        """Extract and return normalised text from a Word document.

        Routes large or legacy files to a child worker subprocess.  For files
        handled in-process, selects the extraction method based on extension
        and logs the method used.

        Raises
        ------
        ValueError
            If *path* has an unsupported extension.
        TextExtractionError
            If all extraction methods fail.
        """
        source = validate_file(path)
        ext = source.suffix.lower().lstrip(".")

        if _should_route_to_office_subprocess(source, ext):
            return run_office_worker(source, config=self._worker_config())

        method_used = ""

        if ext in ("docx", "docm"):
            text, parser_method, paragraphs_count = self._extract_docx(source)
            method_used = parser_method
        elif ext == "doc":
            com_docx = self._try_doc_via_com(source)
            if com_docx is not None:
                try:
                    text, parser_method, paragraphs_count = self._extract_docx(com_docx)
                    method_used = f"word_com+{parser_method}"
                finally:
                    shutil.rmtree(com_docx.parent, ignore_errors=True)
            else:
                with tempfile.TemporaryDirectory(prefix="office_word_") as out_dir:
                    converted = self.converter.convert(
                        source,
                        "docx",
                        Path(out_dir),
                        timeout_s=90,
                    )
                    text, parser_method, paragraphs_count = self._extract_docx(converted)
                    method_used = f"libreoffice+{parser_method}"
        else:
            raise ValueError(f"Unsupported extension for Word extractor: {ext}")

        normalized = _normalize_office_text(text)
        normalized, truncated = _truncate_text(normalized, self.max_output_chars)
        if truncated:
            logger.warning(
                "Word extraction truncated output",
                extra={"source": str(source), "max_output_chars": self.max_output_chars, "truncated": True},
            )

        logger.info(
            "Word extraction completed",
            extra={
                "source": str(source),
                "extractor": "word",
                "method_used": method_used,
                "paragraphs_count": paragraphs_count,
                "truncated": truncated,
                "text_length": len(normalized),
            },
        )
        return normalized

    def _worker_config(self) -> dict:
        """Return the config dict forwarded to the child worker subprocess."""
        return {
            "soffice_path": self.converter.soffice_path,
            "max_output_chars": self.max_output_chars,
        }

    def _try_doc_via_com(self, source: Path) -> Path | None:
        """Convert .doc to a temp .docx using Word COM automation.

        Returns the path to the converted .docx, or None if COM is unavailable
        or conversion fails (caller should fall back to LibreOffice).
        """
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            logger.debug("pywin32 not available; skipping Word COM conversion for .doc")
            return None

        tmp_dir = Path(tempfile.mkdtemp(prefix="office_word_com_"))
        out_path = tmp_dir / f"{source.stem}.docx"

        pythoncom.CoInitialize()
        word = None
        try:
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0
            doc = word.Documents.Open(
                str(source.absolute()),
                ConfirmConversions=False,
                ReadOnly=True,
                AddToRecentFiles=False,
            )
            # wdFormatXMLDocument = 12 (.docx)
            doc.SaveAs2(str(out_path.absolute()), FileFormat=12)
            doc.Close(False)
            logger.info("Word COM converted .doc to .docx: %s -> %s", source, out_path)
            return out_path
        except Exception as exc:
            logger.warning("Word COM conversion failed for %s, will try LibreOffice: %s", source, exc)
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return None
        finally:
            if word is not None:
                try:
                    word.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()

    def _extract_docx(self, path: Path) -> tuple[str, str, int]:
        """Parse a .docx/.docm file and return ``(text, method, paragraph_count)``.

        Tries mammoth first for richer text extraction, falling back to
        python-docx on any mammoth error.  Tables from all methods are
        flattened to tab-separated rows.

        Raises
        ------
        TextExtractionError
            If python-docx is not installed and mammoth also fails.
        """
        mammoth_text = ""
        try:
            import mammoth

            with open(path, "rb") as file_handle:
                result = mammoth.extract_raw_text(file_handle)
            mammoth_text = (result.value or "").strip()
            if mammoth_text:
                paragraphs_count = sum(1 for line in mammoth_text.splitlines() if line.strip())
                return mammoth_text, "mammoth", paragraphs_count
        except Exception as exc:
            logger.warning("Mammoth extraction failed, falling back to python-docx: %s", exc)

        try:
            from docx import Document
        except Exception as exc:
            raise TextExtractionError(
                f"python-docx unavailable for fallback parsing: {path}"
            ) from exc

        doc = Document(str(path))
        paragraphs: list[str] = []
        for paragraph in doc.paragraphs:
            text = (paragraph.text or "").strip()
            if text:
                paragraphs.append(text)

        for table in doc.tables:
            for row in table.rows:
                cells = [(cell.text or "").strip() for cell in row.cells]
                if any(cells):
                    paragraphs.append("\t".join(cells))

        joined = "\n".join(paragraphs)
        return joined, "python-docx", len(paragraphs)


class PresentationTextExtractor(FileTextExtractor):
    """Extract text from PowerPoint presentations (.pptx, .pptm, .ppsx, .ppt, .pps).

    Extraction chain
    ----------------
    * ``.pptx`` / ``.pptm`` / ``.ppsx``: python-pptx (in-process).
    * ``.ppt`` / ``.pps``: Windows COM ``PowerPoint.Application`` → ``.pptx``
      (primary); LibreOffice ``soffice --convert-to pptx`` (fallback).

    Slide notes are included by default.  Shape types handled: text frames,
    tables, and grouped shapes (recursively).
    """

    file_extensions: List[str] = ["pptx", "pptm", "ppsx", "ppt", "pps"]

    def __init__(
        self,
        include_notes: bool = True,
        max_slides: int = 500,
        max_chars_per_slide: int = 50_000,
        converter: OfficeConverter | None = None,
    ):
        """Initialise the extractor.

        Parameters
        ----------
        include_notes:
            When True, speaker notes are appended to each slide block.
        max_slides:
            Slides beyond this count are silently dropped (text is truncated).
        max_chars_per_slide:
            Per-slide character cap; excess characters are discarded.
        converter:
            LibreOffice fallback converter for legacy formats.
        """
        super().__init__()
        self.include_notes = include_notes
        self.max_slides = max_slides
        self.max_chars_per_slide = max_chars_per_slide
        self.converter = converter or OfficeConverter()
    def __call__(self, path: str) -> str:
        """Extract and return normalised text from a presentation file.

        Routes large or legacy files to a child worker subprocess.  For files
        handled in-process, selects the extraction method based on extension
        and logs per-slide metrics.

        Raises
        ------
        ValueError
            If *path* has an unsupported extension.
        TextExtractionError
            If all conversion/extraction methods fail.
        """
        source = validate_file(path)
        ext = source.suffix.lower().lstrip(".")

        if _should_route_to_office_subprocess(source, ext):
            return run_office_worker(source, config=self._worker_config())

        if ext in ("pptx", "pptm", "ppsx"):
            text, slide_count, shapes_scanned, truncated, notes_included = self._extract_pptx(source)
            method_used = "python-pptx"
        elif ext in ("ppt", "pps"):
            com_pptx = self._try_ppt_via_com(source)
            if com_pptx is not None:
                try:
                    text, slide_count, shapes_scanned, truncated, notes_included = self._extract_pptx(com_pptx)
                    method_used = "powerpoint_com+python-pptx"
                finally:
                    shutil.rmtree(com_pptx.parent, ignore_errors=True)
            else:
                with tempfile.TemporaryDirectory(prefix="office_ppt_") as out_dir:
                    converted = self.converter.convert(
                        source,
                        "pptx",
                        Path(out_dir),
                        timeout_s=120,
                    )
                    text, slide_count, shapes_scanned, truncated, notes_included = self._extract_pptx(converted)
                    method_used = "libreoffice+python-pptx"
        else:
            raise ValueError(f"Unsupported extension for presentation extractor: {ext}")

        normalized = _normalize_office_text(text)
        logger.info(
            "Presentation extraction completed",
            extra={
                "source": str(source),
                "extractor": "presentation",
                "method_used": method_used,
                "slide_count": slide_count,
                "shapes_scanned": shapes_scanned,
                "notes_included": notes_included,
                "truncated": truncated,
                "text_length": len(normalized),
            },
        )
        return normalized

    def _worker_config(self) -> dict:
        """Return the config dict forwarded to the child worker subprocess."""
        return {
            "soffice_path": self.converter.soffice_path,
            "include_ppt_notes": self.include_notes,
            "ppt": {
                "max_slides": self.max_slides,
                "max_chars_per_slide": self.max_chars_per_slide,
            },
        }

    def _try_ppt_via_com(self, source: Path) -> Path | None:
        """Convert .ppt/.pps to a temp .pptx using PowerPoint COM automation.

        Returns the path to the converted .pptx, or None if COM is unavailable
        or conversion fails (caller should fall back to LibreOffice).
        """
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            logger.debug("pywin32 not available; skipping PowerPoint COM conversion")
            return None

        tmp_dir = Path(tempfile.mkdtemp(prefix="office_ppt_com_"))
        out_path = tmp_dir / f"{source.stem}.pptx"

        pythoncom.CoInitialize()
        ppt_app = None
        try:
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            # Open without a window
            pres = ppt_app.Presentations.Open(
                str(source.absolute()),
                ReadOnly=True,
                Untitled=True,
                WithWindow=False,
            )
            # ppSaveAsOpenXMLPresentation = 24
            pres.SaveAs(str(out_path.absolute()), FileFormat=24)
            pres.Close()
            logger.info("PowerPoint COM converted to .pptx: %s -> %s", source, out_path)
            return out_path
        except Exception as exc:
            logger.warning("PowerPoint COM conversion failed for %s, will try LibreOffice: %s", source, exc)
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return None
        finally:
            if ppt_app is not None:
                try:
                    ppt_app.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()

    def _extract_pptx(self, path: Path) -> tuple[str, int, int, bool, bool]:
        """Parse a .pptx/.pptm/.ppsx file with python-pptx.

        Returns
        -------
        tuple of (text, slide_count, shapes_scanned, truncated, notes_included)
            *text* is a multi-block string with one ``# slide N`` header per
            slide.  *truncated* is True if ``max_slides`` or
            ``max_chars_per_slide`` was hit.

        Raises
        ------
        TextExtractionError
            If python-pptx is not installed.
        """
        try:
            from pptx import Presentation
        except Exception as exc:
            raise TextExtractionError("python-pptx is required for presentation extraction") from exc

        presentation = Presentation(str(path))
        blocks: list[str] = []
        shapes_scanned = 0
        truncated = False

        for index, slide in enumerate(presentation.slides, start=1):
            if index > self.max_slides:
                truncated = True
                logger.warning(
                    "Presentation truncated by slide limit",
                    extra={"source": str(path), "max_slides": self.max_slides, "truncated": True},
                )
                break

            slide_lines: list[str] = [f"# slide {index}"]
            text_parts: list[str] = []
            for shape in slide.shapes:
                shapes_scanned += 1
                shape_text = self._shape_text(shape)
                if shape_text:
                    text_parts.append(shape_text)

            if text_parts:
                slide_lines.extend(text_parts)

            if self.include_notes:
                notes_text = self._notes_text(slide)
                if notes_text:
                    slide_lines.append("notes:")
                    slide_lines.append(notes_text)

            slide_block = "\n".join(slide_lines).strip()
            if len(slide_block) > self.max_chars_per_slide:
                slide_block = slide_block[: self.max_chars_per_slide]
                truncated = True
                logger.warning(
                    "Presentation slide text truncated",
                    extra={
                        "source": str(path),
                        "slide_index": index,
                        "max_chars_per_slide": self.max_chars_per_slide,
                        "truncated": True,
                    },
                )

            blocks.append(slide_block)

        return "\n\n".join(blocks), min(len(presentation.slides), self.max_slides), shapes_scanned, truncated, self.include_notes

    @staticmethod
    def _shape_text(shape) -> str:
        """Extract all visible text from a single slide shape.

        Handles text frames, tables, and group shapes (recursively).  Returns
        an empty string for shapes with no extractable text (images, charts,
        etc.).
        """
        lines: list[str] = []

        if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                runs = [run.text for run in paragraph.runs if run.text]
                paragraph_text = "".join(runs).strip() if runs else (paragraph.text or "").strip()
                if paragraph_text:
                    lines.append(paragraph_text)

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_cells = [(cell.text or "").strip() for cell in row.cells]
                if any(row_cells):
                    lines.append("\t".join(row_cells))

        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            for grouped_shape in shape.shapes:
                nested_text = PresentationTextExtractor._shape_text(grouped_shape)
                if nested_text:
                    lines.append(nested_text)

        return "\n".join(lines).strip()

    @staticmethod
    def _notes_text(slide) -> str:
        """Return the speaker notes text for *slide*, or an empty string.

        Silently suppresses any error (e.g. slides with no notes slide object)
        so a single malformed notes section does not abort the whole extraction.
        """
        try:
            if not getattr(slide, "has_notes_slide", False):
                return ""
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is None:
                return ""
            return (notes_frame.text or "").strip()
        except Exception:
            return ""


class SpreadsheetTextExtractor(FileTextExtractor):
    """Extract structured text from spreadsheets (.xlsx, .xlsm, .xls).

    Extraction chain
    ----------------
    * ``.xlsx`` / ``.xlsm``: openpyxl (in-process, read-only mode).
    * ``.xls``: Windows COM ``Excel.Application`` → ``.xlsx`` (primary);
      LibreOffice ``soffice --convert-to xlsx`` (fallback).

    Output format
    -------------
    Each sheet is prefixed with a ``# sheet: <name>`` header.  Rows are
    rendered as tab-delimited (or custom *delimiter*) lines.  Trailing empty
    cells are stripped from each row.  Fully empty rows are omitted.
    """

    file_extensions: List[str] = ["xlsx", "xlsm", "xls"]

    def __init__(
        self,
        delimiter: str = "\t",
        max_sheets: int = 20,
        max_rows_per_sheet: int = 5_000,
        max_cols_per_sheet: int = 200,
        max_total_cells: int = 500_000,
        converter: OfficeConverter | None = None,
    ):
        """Initialise the extractor.

        Parameters
        ----------
        delimiter:
            Column separator used when serialising rows to text.
        max_sheets:
            Sheets beyond this count are silently dropped.
        max_rows_per_sheet:
            Rows beyond this count within a sheet are silently dropped.
        max_cols_per_sheet:
            Columns beyond this count within a sheet are silently dropped.
        max_total_cells:
            Hard cap across all sheets; extraction stops when reached.
        converter:
            LibreOffice fallback converter for ``.xls`` files.
        """
        super().__init__()
        self.delimiter = delimiter
        self.max_sheets = max_sheets
        self.max_rows_per_sheet = max_rows_per_sheet
        self.max_cols_per_sheet = max_cols_per_sheet
        self.max_total_cells = max_total_cells
        self.converter = converter or OfficeConverter()
    def __call__(self, path: str) -> str:
        """Extract and return normalised text from a spreadsheet file.

        Routes large or legacy files to a child worker subprocess.  For files
        handled in-process, selects the extraction method based on extension
        and logs sheet/row/cell metrics.

        Raises
        ------
        ValueError
            If *path* has an unsupported extension.
        TextExtractionError
            If all conversion/extraction methods fail.
        """
        source = validate_file(path)
        ext = source.suffix.lower().lstrip(".")

        if _should_route_to_office_subprocess(source, ext):
            return run_office_worker(source, config=self._worker_config())

        if ext in ("xlsx", "xlsm"):
            text, sheets_scanned, rows_scanned, cells_scanned, truncated = self._extract_xlsx(source)
            method_used = "openpyxl"
        elif ext == "xls":
            com_xlsx = self._try_xls_via_com(source)
            if com_xlsx is not None:
                try:
                    text, sheets_scanned, rows_scanned, cells_scanned, truncated = self._extract_xlsx(com_xlsx)
                    method_used = "excel_com+openpyxl"
                finally:
                    shutil.rmtree(com_xlsx.parent, ignore_errors=True)
            else:
                with tempfile.TemporaryDirectory(prefix="office_xls_") as out_dir:
                    converted = self.converter.convert(
                        source,
                        "xlsx",
                        Path(out_dir),
                        timeout_s=90,
                    )
                    text, sheets_scanned, rows_scanned, cells_scanned, truncated = self._extract_xlsx(converted)
                    method_used = "libreoffice+openpyxl"
        else:
            raise ValueError(f"Unsupported extension for spreadsheet extractor: {ext}")

        normalized = _normalize_office_text(text)
        logger.info(
            "Spreadsheet extraction completed",
            extra={
                "source": str(source),
                "extractor": "spreadsheet",
                "method_used": method_used,
                "sheets_scanned": sheets_scanned,
                "rows_scanned": rows_scanned,
                "cells_scanned": cells_scanned,
                "truncated": truncated,
                "text_length": len(normalized),
            },
        )
        return normalized

    def _worker_config(self) -> dict:
        """Return the config dict forwarded to the child worker subprocess."""
        return {
            "soffice_path": self.converter.soffice_path,
            "xlsx": {
                "max_sheets": self.max_sheets,
                "max_rows_per_sheet": self.max_rows_per_sheet,
                "max_cols_per_sheet": self.max_cols_per_sheet,
                "max_total_cells": self.max_total_cells,
            },
        }

    def _try_xls_via_com(self, source: Path) -> Path | None:
        """Convert .xls to a temp .xlsx using Excel COM automation.

        Returns the path to the converted .xlsx, or None if COM is unavailable
        or conversion fails (caller should fall back to LibreOffice).
        """
        try:
            import pythoncom
            import win32com.client
        except ImportError:
            logger.debug("pywin32 not available; skipping Excel COM conversion for .xls")
            return None

        tmp_dir = Path(tempfile.mkdtemp(prefix="office_xls_com_"))
        out_path = tmp_dir / f"{source.stem}.xlsx"

        pythoncom.CoInitialize()
        excel_app = None
        try:
            excel_app = win32com.client.Dispatch("Excel.Application")
            excel_app.Visible = False
            excel_app.DisplayAlerts = False
            wb = excel_app.Workbooks.Open(
                str(source.absolute()),
                UpdateLinks=False,
                ReadOnly=True,
                AddToMru=False,
            )
            # xlOpenXMLWorkbook = 51 (.xlsx)
            wb.SaveAs(str(out_path.absolute()), FileFormat=51)
            wb.Close(False)
            logger.info("Excel COM converted .xls to .xlsx: %s -> %s", source, out_path)
            return out_path
        except Exception as exc:
            logger.warning("Excel COM conversion failed for %s, will try LibreOffice: %s", source, exc)
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return None
        finally:
            if excel_app is not None:
                try:
                    excel_app.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()

    def _extract_xlsx(self, path: Path) -> tuple[str, int, int, int, bool]:
        """Parse a .xlsx/.xlsm file with openpyxl in read-only / data-only mode.

        Returns
        -------
        tuple of (text, sheets_scanned, rows_scanned, cells_scanned, truncated)
            *text* is a multi-block string, one block per sheet.  *truncated*
            is True if any limit (sheets, rows, total cells) was hit.

        Raises
        ------
        TextExtractionError
            If openpyxl is not installed.
        """
        try:
            from openpyxl import load_workbook
        except Exception as exc:
            raise TextExtractionError("openpyxl is required for spreadsheet extraction") from exc

        workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
        blocks: list[str] = []
        rows_scanned = 0
        cells_scanned = 0
        sheets_scanned = 0
        truncated = False

        try:
            for sheet_name in workbook.sheetnames:
                if sheets_scanned >= self.max_sheets:
                    truncated = True
                    logger.warning(
                        "Spreadsheet truncated by sheet limit",
                        extra={"source": str(path), "max_sheets": self.max_sheets, "truncated": True},
                    )
                    break

                worksheet = workbook[sheet_name]
                sheet_lines = [f"# sheet: {sheet_name}"]
                sheet_rows = 0

                for row in worksheet.iter_rows(
                    min_row=1,
                    max_row=self.max_rows_per_sheet,
                    min_col=1,
                    max_col=self.max_cols_per_sheet,
                    values_only=True,
                ):
                    row_cells: list[str] = []
                    for cell in row:
                        if cells_scanned >= self.max_total_cells:
                            truncated = True
                            break
                        row_cells.append(self._cell_to_text(cell))
                        cells_scanned += 1

                    if truncated:
                        logger.warning(
                            "Spreadsheet truncated by total cell limit",
                            extra={
                                "source": str(path),
                                "max_total_cells": self.max_total_cells,
                                "truncated": True,
                            },
                        )
                        break

                    if any(value for value in row_cells):
                        while row_cells and row_cells[-1] == "":
                            row_cells.pop()
                        sheet_lines.append(self.delimiter.join(row_cells))
                    rows_scanned += 1
                    sheet_rows += 1

                if sheet_rows >= self.max_rows_per_sheet and worksheet.max_row and worksheet.max_row > self.max_rows_per_sheet:
                    truncated = True
                    logger.warning(
                        "Spreadsheet sheet truncated by row limit",
                        extra={
                            "source": str(path),
                            "sheet": sheet_name,
                            "max_rows_per_sheet": self.max_rows_per_sheet,
                            "truncated": True,
                        },
                    )

                blocks.append("\n".join(sheet_lines))
                sheets_scanned += 1

                if truncated and cells_scanned >= self.max_total_cells:
                    break
        finally:
            workbook.close()

        return "\n\n".join(blocks), sheets_scanned, rows_scanned, cells_scanned, truncated

    @staticmethod
    def _cell_to_text(value) -> str:
        """Convert a cell value to a stripped string, returning "" for None."""
        if value is None:
            return ""
        return str(value).strip()
